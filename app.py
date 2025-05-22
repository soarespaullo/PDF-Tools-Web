from flask import Flask, render_template, request, send_file, redirect, flash, url_for
from werkzeug.utils import secure_filename
import os
import zipfile
import PyPDF2
import subprocess
import pytesseract
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PyPDF2.errors import FileNotDecryptedError, PdfReadError
from weasyprint import HTML
from datetime import datetime
from PIL import Image, UnidentifiedImageError
from fpdf import FPDF
from pdf2image import convert_from_path
from pdf2docx import Converter
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from io import BytesIO
from bs4 import BeautifulSoup
from bleach import clean

# Cria a aplicação Flask
app = Flask(__name__)

# Define uma chave secreta para a aplicação
# Ela é usada para manter dados seguros nas sessões (por exemplo, cookies)
app.secret_key = 'supersecretkey'

# Define o caminho absoluto do diretório onde o arquivo atual (geralmente app.py) está localizado
BASE_DIR = os.path.abspath(os.path.dirname(__file__))

# Cria o caminho completo para a pasta 'uploads', usada para armazenar arquivos enviados temporariamente
RESULT_FOLDER = os.path.join(BASE_DIR, 'results')

# Cria o caminho completo para a pasta 'results', usada para armazenar arquivos gerados pelo app
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')

# Garante que a pasta 'uploads' exista; se não existir, será criada
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Garante que a pasta 'results' exista; se não existir, será criada
os.makedirs(RESULT_FOLDER, exist_ok=True)

# --- ROTA: Exibir a Página Inicial com os Formulários ---
@app.route('/')
def index():
    return render_template('index.html')

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: Juntar PDFs ---
@app.route('/juntar', methods=['POST'])
def juntar_pdf():
    # Obtém todos os arquivos enviados no campo "pdfs" (com multiple no HTML)
    arquivos = request.files.getlist("pdfs")

    # Cria o objeto responsável por juntar os PDFs
    merger = PyPDF2.PdfMerger()

    # Contador de PDFs válidos
    pdfs_validos = 0

    # Lista para armazenar os nomes dos arquivos inválidos
    arquivos_invalidos = []

    # Percorre cada arquivo enviado
    for file in arquivos:
        # Verifica se o arquivo existe e tem a extensão .pdf
        if file and file.filename.lower().endswith('.pdf'):
            # Garante um nome de arquivo seguro
            filename = secure_filename(file.filename)
            # Define o caminho temporário onde o arquivo será salvo
            path = os.path.join(UPLOAD_FOLDER, filename)
            # Salva o arquivo no servidor temporariamente
            file.save(path)

            try:
                # Tenta abrir o arquivo com PyPDF2 para validar se é um PDF de verdade
                with open(path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    _ = len(reader.pages)  # Verifica se possui páginas

                # Se passou pela validação, adiciona ao merger
                merger.append(path)
                pdfs_validos += 1

            except Exception:
                # Remove o arquivo corrompido ou inválido
                os.remove(path)
                # Adiciona à lista de arquivos inválidos para exibir mensagem depois
                arquivos_invalidos.append(filename)

        else:
            # Arquivos que não têm extensão .pdf
            arquivos_invalidos.append(file.filename)

    # Exibe uma mensagem flash para cada arquivo inválido detectado
    for nome in arquivos_invalidos:
        flash(f"❌ O arquivo '{nome}' não é um PDF válido e foi ignorado.")

    # Se menos de dois arquivos válidos foram enviados,
    # e não houve arquivos inválidos (já tratados com mensagens),
    # então mostra mensagem genérica de erro
    if pdfs_validos < 2:
        if not arquivos_invalidos:
            flash("⚠️ Envie pelo menos dois arquivos PDF válidos para juntar.")
        return redirect(request.referrer or url_for('index'))

    # Cria um nome único para o PDF de saída com base na data/hora
    data_formatada = datetime.now().strftime('%d-%m-%Y_%H-%M-%S')
    output_filename = f"juntado_{data_formatada}.pdf"
    output_path = os.path.join(RESULT_FOLDER, output_filename)

    # Escreve o PDF final juntado
    merger.write(output_path)
    merger.close()

    # Exibe mensagem de sucesso
    flash("✅ PDFs juntados com sucesso!")

    # Redireciona para a página que mostrará o link para baixar
    return redirect(url_for('download_page', filename=output_filename))


# Rota que exibe a página HTML com o link de download do arquivo gerado
@app.route('/download/<filename>')
def download_page(filename):
    # Renderiza o template 'download.html' passando o nome do arquivo como variável
    return render_template('download.html', filename=filename)


# Rota que permite baixar diretamente o arquivo PDF gerado
@app.route('/baixar/<filename>')
def baixar_pdf(filename):
    # Monta o caminho completo até o arquivo no servidor
    path = os.path.join(RESULT_FOLDER, filename)
    # Envia o arquivo como anexo (força o download no navegador)
    return send_file(path, as_attachment=True)

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: Dividir PDF ---
@app.route('/dividir', methods=['POST'])
def dividir_pdf():
    # Verifica se o arquivo foi enviado no formulário
    if 'pdf' not in request.files:
        flash("⚠️ Nenhum arquivo foi enviado.")
        return redirect(request.referrer or url_for('index'))

    file = request.files['pdf']

    # Verifica se o nome do arquivo está vazio
    if file.filename == '':
        flash("⚠️ Nenhum arquivo selecionado.")
        return redirect(request.referrer or url_for('index'))

    # Verifica se a extensão é .pdf
    if not file.filename.lower().endswith('.pdf'):
        flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
        return redirect(request.referrer or url_for('index'))

    # Garante um nome seguro e salva o arquivo
    filename = secure_filename(file.filename)
    path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(path)

    try:
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)

            # Verifica se está criptografado
            if reader.is_encrypted:
                os.remove(path)
                flash("🔒 O arquivo PDF está protegido por senha e não pode ser processado.")
                return redirect(request.referrer or url_for('index'))

            # Verifica se há páginas no PDF
            if len(reader.pages) == 0:
                os.remove(path)
                flash("⚠️ O PDF está vazio.")
                return redirect(request.referrer or url_for('index'))

            # Define nome e caminho do ZIP final
            zip_filename = f"{os.path.splitext(filename)[0]}_paginas.zip"
            zip_path = os.path.join(RESULT_FOLDER, zip_filename)
            paginas_geradas = []

            # Cria o ZIP com os PDFs de cada página
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for i, page in enumerate(reader.pages):
                    writer = PyPDF2.PdfWriter()
                    writer.add_page(page)

                    output_filename = f"{os.path.splitext(filename)[0]}_pagina_{i+1}.pdf"
                    output_path = os.path.join(RESULT_FOLDER, output_filename)

                    with open(output_path, 'wb') as out:
                        writer.write(out)

                    zipf.write(output_path, arcname=output_filename)
                    paginas_geradas.append(output_filename)

    except PyPDF2.errors.PdfReadError:
        os.remove(path)
        flash(f"❌ O arquivo '{file.filename}' não é um PDF válido ou está corrompido.")
        return redirect(request.referrer or url_for('index'))

    except Exception:
        os.remove(path)
        flash(f"❌ Ocorreu um erro ao processar o arquivo '{file.filename}'. Certifique-se de que seja um PDF válido.")
        return redirect(request.referrer or url_for('index'))

    # Mensagem de sucesso
    flash("✅ PDF dividido com sucesso!")

    # Redireciona para página com links de download
    return redirect(url_for('download_paginas',
                            arquivos=','.join(paginas_geradas),
                            zipfile=zip_filename))

# Rota para exibir a página de download das páginas geradas
@app.route('/paginas-divididas')
def download_paginas():
    # Recupera os nomes dos arquivos individuais e do ZIP dos parâmetros da URL
    arquivos = request.args.get('arquivos', '').split(',')
    zipfile_name = request.args.get('zipfile', '')

    # Renderiza o template com os nomes dos arquivos
    return render_template('paginas_divididas.html', arquivos=arquivos, zipfile=zipfile_name)

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: Comprimir PDF ---
@app.route('/comprimir', methods=['POST'])
def comprimir_pdf():
    # Verifica se o arquivo foi enviado no formulário
    if 'pdf' not in request.files:
        flash("⚠️ Nenhum arquivo foi enviado.")
        return redirect(request.referrer or url_for('index'))

    file = request.files['pdf']

    # Verifica se o nome do arquivo está vazio
    if file.filename == '':
        flash("⚠️ Nenhum arquivo selecionado.")
        return redirect(request.referrer or url_for('index'))

    # Verifica se a extensão do arquivo é .pdf (não garante que é PDF real, apenas a extensão)
    if not file.filename.lower().endswith('.pdf'):
        flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
        return redirect(request.referrer or url_for('index'))

    # Garante um nome seguro e define o caminho onde o arquivo será salvo
    filename = secure_filename(file.filename)
    input_path = os.path.join(UPLOAD_FOLDER, filename)

    # Salva o arquivo enviado no diretório de uploads
    file.save(input_path)

    # Tenta validar se o arquivo é realmente um PDF e se não está criptografado
    try:
        with open(input_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)

            # Se o PDF estiver criptografado, exibe mensagem apropriada
            if reader.is_encrypted:
                raise Exception("🔒 O arquivo PDF está protegido por senha.")

            # Se o PDF não tiver páginas, também é considerado inválido
            if len(reader.pages) == 0:
                raise Exception("O PDF está vazio.")

    except Exception as e:
        # Remove o arquivo inválido
        os.remove(input_path)

        # Exibe mensagem apropriada dependendo do tipo de erro detectado
        if "protegido" in str(e).lower():
            flash("🔒 O arquivo PDF está protegido por senha.")
        else:
            flash("❌ O arquivo enviado está corrompido ou não é um PDF válido.")
        return redirect(request.referrer or url_for('index'))

    # Define o nome final do arquivo comprimido
    compressed_filename = f"comprimido_{filename}"
    output_path = os.path.join(RESULT_FOLDER, compressed_filename)

    # Comando Ghostscript para comprimir o PDF
    gs_command = [
        "gs",                         # Chama o Ghostscript
        "-sDEVICE=pdfwrite",          # Define o tipo de saída como PDF
        "-dCompatibilityLevel=1.4",   # Define a compatibilidade do PDF
        "-dPDFSETTINGS=/screen",      # Nível de compressão: screen = alta compressão, baixa qualidade
        "-dNOPAUSE",                  # Evita pausas entre páginas
        "-dQUIET",                    # Suprime mensagens de log no terminal
        "-dBATCH",                    # Finaliza o processo automaticamente
        f"-sOutputFile={output_path}",# Arquivo de saída comprimido
        input_path                    # Arquivo original (de entrada)
    ]

    try:
        # Executa o comando do Ghostscript para comprimir o PDF
        subprocess.run(gs_command, check=True)

        # Se sucesso, exibe mensagem positiva ao usuário
        flash("✅ PDF comprimido com sucesso!")

        # Redireciona para a página com o link de download do arquivo comprimido
        return render_template('comprimido.html', filename=compressed_filename)

    except subprocess.CalledProcessError:
        # Caso o Ghostscript não esteja instalado ou ocorra erro na execução
        flash("❌ Erro ao comprimir o PDF. Verifique se o Ghostscript está instalado.")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: Proteger PDF ---
@app.route('/proteger', methods=['POST'])
def proteger_pdf():
    # Verifica se o arquivo foi enviado no formulário
    if 'pdf' not in request.files:
        flash("⚠️ Nenhum arquivo foi enviado.")
        return redirect(request.referrer or url_for('index'))

    file = request.files['pdf']
    senha = request.form.get('senha', '')

    # Verifica se o nome do arquivo está vazio
    if file.filename == '':
        flash("⚠️ Nenhum arquivo selecionado.")
        return redirect(request.referrer or url_for('index'))

    # Verifica se a extensão é .pdf
    if not file.filename.lower().endswith('.pdf'):
        flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
        return redirect(request.referrer or url_for('index'))

    # Verifica se a senha foi fornecida
    if not senha.strip():
        flash("⚠️ Você precisa fornecer uma senha para proteger o PDF.")
        return redirect(request.referrer or url_for('index'))

    # Define nomes e caminhos de entrada/saída
    filename = secure_filename(file.filename)
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_filename = f"protegido_{filename}"
    output_path = os.path.join(RESULT_FOLDER, output_filename)

    # Salva o arquivo no diretório de uploads temporários
    file.save(input_path)

    try:
        # Tenta abrir e validar se é realmente um PDF
        with open(input_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)

            # Verifica se o PDF está criptografado
            if reader.is_encrypted:
                flash("⚠️ Este PDF já está protegido por senha.")
                return redirect(request.referrer or url_for('index'))

            # Verifica se o PDF possui páginas
            if len(reader.pages) == 0:
                raise Exception("O PDF está vazio.")

            # Cria novo writer e adiciona as páginas
            writer = PyPDF2.PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            # Aplica a senha ao PDF de saída
            writer.encrypt(senha)

            # Salva o novo PDF protegido
            with open(output_path, 'wb') as out:
                writer.write(out)

            # Mensagem de sucesso
            flash("🔒 PDF protegido com sucesso!")
            return render_template('protegido.html', filename=output_filename)

    except Exception as e:
        # Remove o arquivo inválido
        os.remove(input_path)

        # Personaliza a mensagem de erro de acordo com o tipo
        if "protegido" in str(e).lower():
            flash("⚠️ Este PDF já está protegido por senha.")
        elif "eof" in str(e).lower() or "startxref" in str(e).lower():
            flash("❌ O arquivo enviado não é um PDF válido ou está corrompido.")
        else:
            flash("❌ Ocorreu um erro ao processar o PDF. Verifique se é um arquivo válido.")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: Desbloquear PDF ---
@app.route('/desbloquear', methods=['POST'])
def desbloquear_pdf():
    # Verifica se o arquivo foi enviado no formulário
    if 'pdf' not in request.files:
        flash("⚠️ Nenhum arquivo foi enviado.")
        return redirect(request.referrer or url_for('index'))

    file = request.files['pdf']
    senha = request.form.get('senha', '')

    # Verifica se o nome do arquivo está vazio
    if file.filename == '':
        flash("⚠️ Nenhum arquivo selecionado.")
        return redirect(request.referrer or url_for('index'))

    # Verifica se a extensão é .pdf
    if not file.filename.lower().endswith('.pdf'):
        flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
        return redirect(request.referrer or url_for('index'))

    # Verifica se a senha foi fornecida
    if not senha.strip():
        flash("⚠️ Você precisa fornecer a senha para desbloquear o PDF.")
        return redirect(request.referrer or url_for('index'))

    # Define nomes e caminhos de entrada/saída
    filename = secure_filename(file.filename)
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_filename = f"desbloqueado_{filename}"
    output_path = os.path.join(RESULT_FOLDER, output_filename)

    # Salva o arquivo temporariamente
    file.save(input_path)

    try:
        # Cria o leitor do PDF
        reader = PyPDF2.PdfReader(input_path)

        # Verifica se o PDF está criptografado
        if reader.is_encrypted:
            try:
                # Tenta desbloquear com a senha fornecida
                if not reader.decrypt(senha):
                    flash("❌ Senha incorreta ou falha ao desbloquear o PDF.")
                    return redirect(request.referrer or url_for('index'))

                # Força acesso a uma página para verificar se desbloqueou mesmo
                _ = reader.pages[0]

            except PyPDF2.errors.FileNotDecryptedError:
                flash("❌ Senha incorreta. O PDF não pôde ser desbloqueado.")
                return redirect(request.referrer or url_for('index'))

            except Exception as e:
                flash("❌ Erro ao processar o PDF. Verifique se o arquivo está corrompido ou é um PDF válido.")
                return redirect(request.referrer or url_for('index'))

            # Se o desbloqueio foi bem-sucedido, cria novo PDF sem senha
            writer = PyPDF2.PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            # Salva o novo PDF desbloqueado
            with open(output_path, 'wb') as f:
                writer.write(f)

            flash("✅ PDF desbloqueado com sucesso!")
            return render_template('desbloqueado.html', filename=output_filename)

        else:
            # O PDF já está desbloqueado
            flash("ℹ️ Este PDF não está protegido por senha.")
            return redirect(request.referrer or url_for('index'))

    except Exception as e:
        # Remove o arquivo temporário se falhar
        os.remove(input_path)
        if "eof" in str(e).lower() or "startxref" in str(e).lower():
            flash("❌ O arquivo enviado não é um PDF válido ou está corrompido.")
        else:
            flash("❌ Ocorreu um erro ao processar o PDF.")
        return redirect(request.referrer or url_for('index'))
        
# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: HTML para PDF ---
@app.route('/html_para_pdf', methods=['POST'])
def html_para_pdf():
    try:
        html_content = request.form.get('html', '')

        if not html_content.strip():
            flash("⚠️ O campo de conteúdo HTML está vazio.")
            return redirect(request.referrer or url_for('index'))

        # Sanitiza o conteúdo HTML
        safe_html = clean(
            html_content,
            tags=[
                'p', 'br', 'b', 'i', 'u', 'strong', 'em', 'h1', 'h2', 'h3', 'h4',
                'ul', 'ol', 'li', 'table', 'thead', 'tbody', 'tr', 'td', 'th',
                'a', 'img', 'div', 'span', 'hr'
            ],
            attributes={
                'a': ['href', 'title'],
                'img': ['src', 'alt', 'width', 'height'],
                '*': ['style']
            },
            protocols=['http', 'https', 'mailto'],
            strip=True
        )

        # Valida se há pelo menos uma tag HTML após sanitização
        soup = BeautifulSoup(safe_html, 'html.parser')
        if not soup.find():  # Nenhuma tag HTML encontrada
            flash("⚠️ O conteúdo enviado não parece ser um HTML válido.")
            return redirect(request.referrer or url_for('index'))

        # Define nome e caminho do arquivo PDF
        timestamp = datetime.now().strftime("%d-%m-%Y_%H-%M-%S")
        output_filename = f"html_convertido_{timestamp}.pdf"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Gera o PDF com WeasyPrint
        HTML(string=safe_html).write_pdf(output_path)

        flash("✅ HTML convertido para PDF com sucesso!")
        return render_template('download_unico.html', filename=output_filename)

    except Exception:
        flash("❌ Ocorreu um erro ao converter o HTML. Verifique se o conteúdo está bem formatado.")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: PDF para JPG ---
@app.route('/pdf_para_jpg', methods=['POST'])
def pdf_para_jpg():
    try:
        # Verifica se o arquivo foi enviado via formulário
        if 'pdf' not in request.files:
            flash("⚠️ Nenhum arquivo foi enviado.")
            return redirect(request.referrer or url_for('index'))

        file = request.files['pdf']

        # Verifica se o usuário selecionou algum arquivo
        if file.filename == '':
            flash("⚠️ Nenhum arquivo selecionado.")
            return redirect(request.referrer or url_for('index'))

        # Verifica se o arquivo tem extensão .pdf
        if not file.filename.lower().endswith('.pdf'):
            flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
            return redirect(request.referrer or url_for('index'))

        # Gera timestamp no formato dd-mm-aaaa_HH-MM-SS
        timestamp = datetime.now().strftime("%d-%m-%Y_%H-%M-%S")
        safe_name = secure_filename(file.filename)
        filename = f"{timestamp}_{safe_name}"

        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        # Verifica se o PDF está protegido por senha
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    os.remove(input_path)
                    flash("🔒 O arquivo PDF está protegido por senha.")
                    return redirect(request.referrer or url_for('index'))
        except PdfReadError as e:
            # Trata erro ao ler o PDF (ex: PDF corrompido ou truncado)
            os.remove(input_path)
            if "EOF marker not found" in str(e):
                flash("❌ O arquivo PDF parece estar incompleto ou corrompido.")
            else:
                flash(f"❌ Erro ao ler o PDF: {str(e)}")
            return redirect(request.referrer or url_for('index'))

        # Converte cada página do PDF em uma imagem JPG
        imagens = convert_from_path(input_path)
        arquivos_gerados = []

        for i, img in enumerate(imagens):
            output_filename = f"{filename.rsplit('.', 1)[0]}_pagina_{i+1}.jpg"
            output_path = os.path.join(RESULT_FOLDER, output_filename)
            img.save(output_path, 'JPEG')
            arquivos_gerados.append(output_filename)

        # Mensagem de sucesso + renderiza página com os arquivos gerados
        flash("✅ Conversão para JPG concluída!")
        return render_template('paginas_convertidas.html', arquivos=arquivos_gerados)

    except Exception as e:
        # Remove o arquivo temporário em caso de erro
        if 'input_path' in locals() and os.path.exists(input_path):
            os.remove(input_path)
        # Exibe erro genérico para falhas inesperadas
        flash(f"❌ Erro ao converter PDF para JPG: {str(e)}")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: JPG para PDF ---
@app.route('/jpg_para_pdf', methods=['POST'])
def jpg_para_pdf():
    try:
        imagens = request.files.getlist("imagens")

        # Verifica se alguma imagem foi selecionada
        if not imagens or imagens[0].filename == '':
            flash("❌ Nenhuma imagem selecionada para conversão.")
            return redirect('/')

        pdf = FPDF()
        pdf.set_auto_page_break(0)  # Evita quebra automática de página

        for img in imagens:
            filename = secure_filename(img.filename)

            # Antes de salvar, valida se o arquivo é realmente uma imagem válida
            try:
                # Tenta abrir o arquivo com PIL para validar
                image = Image.open(img)
                image.verify()  # Verifica se a imagem está íntegra e é válida
                img.seek(0)     # Reset para início do arquivo após verify (importante)
            except (UnidentifiedImageError, IOError):
                flash(f"❌ O arquivo '{filename}' não é uma imagem válida.")
                return redirect('/')

            # Salva o arquivo validado na pasta temporária
            path = os.path.join(UPLOAD_FOLDER, filename)
            img.save(path)

            # Reabre a imagem para converter para RGB (FPDF só aceita JPEG RGB)
            image = Image.open(path).convert('RGB')

            # Salva temporariamente em JPEG para compatibilidade com FPDF
            image_jpg_path = path.rsplit('.', 1)[0] + ".jpg"
            image.save(image_jpg_path)

            # Adiciona página no PDF e insere a imagem JPEG
            pdf.add_page()
            pdf.image(image_jpg_path, x=10, y=10, w=190)

        # Gera um nome amigável para o PDF resultante com timestamp
        output_filename = f"convertido_{datetime.now().strftime('%d-%m-%Y_%H-%M-%S')}.pdf"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Salva o PDF final
        pdf.output(output_path)

        flash("✅ Imagens convertidas para PDF com sucesso!")
        return render_template('download_unico.html', filename=output_filename)

    except Exception as e:
        flash(f"❌ Erro ao converter JPG para PDF: {str(e)}")
        return redirect('/')

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: OCR PDF ---
@app.route('/ocr_pdf', methods=['POST'])
def ocr_pdf():
    try:
        # Verifica se o arquivo foi enviado no formulário
        if 'pdf' not in request.files:
            flash("⚠️ Nenhum arquivo foi enviado.")
            return redirect(request.referrer or url_for('index'))

        file = request.files['pdf']

        # Verifica se o usuário selecionou um arquivo
        if file.filename == '':
            flash("⚠️ Nenhum arquivo selecionado.")
            return redirect(request.referrer or url_for('index'))

        # Verifica se o arquivo tem extensão .pdf (simples validação)
        if not file.filename.lower().endswith('.pdf'):
            flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
            return redirect(request.referrer or url_for('index'))

        # Garante um nome seguro e salva o arquivo na pasta temporária
        filename = secure_filename(file.filename)
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        # Tenta abrir o PDF para verificar se está válido e não protegido
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    os.remove(input_path)
                    flash("🔒 O arquivo PDF está protegido por senha.")
                    return redirect(request.referrer or url_for('index'))
        except PdfReadError as e:
            # Trata erro de PDF inválido ou corrompido
            os.remove(input_path)
            flash("❌ O arquivo PDF parece estar incompleto ou corrompido.")
            return redirect(request.referrer or url_for('index'))

        # Converte o PDF em imagens para realizar OCR
        imagens = convert_from_path(input_path)

        # Inicializa string para juntar o texto extraído
        texto_extraido = ""
        for img in imagens:
            # Extrai texto de cada página imagem usando pytesseract
            texto_extraido += pytesseract.image_to_string(img)

        # Define nome de saída do arquivo texto com timestamp legível
        output_filename = f"ocr_{datetime.now().strftime('%d-%m-%Y_%H-%M-%S')}.txt"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Salva o texto extraído em arquivo UTF-8
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(texto_extraido)

        flash("✅ OCR realizado com sucesso! Texto extraído do PDF.")
        return render_template('download_unico.html', filename=output_filename)

    except Exception as e:
        # Em caso de erro, tenta remover arquivo temporário e mostra mensagem
        if 'input_path' in locals() and os.path.exists(input_path):
            os.remove(input_path)
        flash(f"❌ Erro ao realizar OCR: {str(e)}")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: PDF para Word ---
@app.route('/pdf_para_word', methods=['POST'])
def pdf_para_word():
    try:
        # Verifica se o arquivo foi enviado no formulário
        if 'pdf' not in request.files:
            flash("⚠️ Nenhum arquivo foi enviado.")
            return redirect(request.referrer or url_for('index'))

        file = request.files['pdf']

        # Verifica se o usuário selecionou algum arquivo
        if file.filename == '':
            flash("⚠️ Nenhum arquivo selecionado.")
            return redirect(request.referrer or url_for('index'))

        # Verifica se o arquivo possui extensão .pdf (validação simples)
        if not file.filename.lower().endswith('.pdf'):
            flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
            return redirect(request.referrer or url_for('index'))

        # Garante nome seguro e salva arquivo na pasta temporária
        filename = secure_filename(file.filename)
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        # Tenta abrir o PDF para garantir que está válido e não protegido
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    os.remove(input_path)
                    flash("🔒 O arquivo PDF está protegido por senha.")
                    return redirect(request.referrer or url_for('index'))
        except PdfReadError:
            # PDF inválido, corrompido ou incompleto
            os.remove(input_path)
            flash("❌ O arquivo PDF parece estar incompleto ou corrompido.")
            return redirect(request.referrer or url_for('index'))

        # Define nome e caminho do arquivo Word de saída com timestamp legível
        output_filename = f"convertido_{datetime.now().strftime('%d-%m-%Y_%H-%M-%S')}.docx"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Converte o PDF para Word usando pdf2docx
        converter = Converter(input_path)
        converter.convert(output_path, start=0, end=None)
        converter.close()

        # Mensagem de sucesso e renderiza página para download
        flash("✅ PDF convertido para Word com sucesso!")
        return render_template('download_unico.html', filename=output_filename)

    except Exception as e:
        # Em caso de erro, tenta remover o arquivo temporário e exibe mensagem
        if 'input_path' in locals() and os.path.exists(input_path):
            os.remove(input_path)
        flash(f"❌ Erro ao converter PDF para Word: {str(e)}")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: PDF para Excel ---
@app.route('/pdf_para_excel', methods=['POST'])
def pdf_para_excel():
    try:
        # Verifica se o arquivo foi enviado no formulário
        if 'pdf' not in request.files:
            flash("⚠️ Nenhum arquivo foi enviado.")
            return redirect(request.referrer or url_for('index'))

        file = request.files['pdf']

        # Verifica se o usuário selecionou algum arquivo
        if file.filename == '':
            flash("⚠️ Nenhum arquivo selecionado.")
            return redirect(request.referrer or url_for('index'))

        # Verifica se o arquivo possui extensão .pdf (validação básica)
        if not file.filename.lower().endswith('.pdf'):
            flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
            return redirect(request.referrer or url_for('index'))

        # Garante nome seguro e salva o PDF na pasta temporária
        filename = secure_filename(file.filename)
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        # Tenta abrir o PDF para garantir que não está protegido nem corrompido
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    os.remove(input_path)
                    flash("🔒 O arquivo PDF está protegido por senha.")
                    return redirect(request.referrer or url_for('index'))
        except PdfReadError:
            # PDF inválido, corrompido ou incompleto
            os.remove(input_path)
            flash("❌ O arquivo PDF parece estar incompleto ou corrompido.")
            return redirect(request.referrer or url_for('index'))

        # Define nome do Excel de saída com timestamp legível
        output_filename = f"tabelas_extraidas_{datetime.now().strftime('%d-%m-%Y_%H-%M-%S')}.xlsx"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Extrai tabelas do PDF usando pdfplumber
        todas_tabelas = []
        with pdfplumber.open(input_path) as pdf:
            for pagina in pdf.pages:
                tabelas = pagina.extract_tables()
                for tabela in tabelas:
                    # Cria DataFrame para cada tabela encontrada (pulando cabeçalho na primeira linha)
                    df = pd.DataFrame(tabela[1:], columns=tabela[0])
                    todas_tabelas.append(df)

        # Se nenhuma tabela for encontrada, avisa o usuário e retorna
        if not todas_tabelas:
            os.remove(input_path)
            flash("⚠️ Nenhuma tabela foi encontrada no PDF.")
            return redirect(request.referrer or url_for('index'))

        # Salva todas as tabelas extraídas em planilhas separadas no arquivo Excel
        with pd.ExcelWriter(output_path) as writer:
            for i, tabela in enumerate(todas_tabelas):
                tabela.to_excel(writer, sheet_name=f"Tabela_{i+1}", index=False)

        # Remove o PDF temporário
        os.remove(input_path)

        # Mensagem de sucesso e renderiza a página de download
        flash("✅ PDF convertido para Excel com sucesso!")
        return render_template('download_unico.html', filename=output_filename)

    except Exception as e:
        # Remove o arquivo temporário em caso de erro
        if 'input_path' in locals() and os.path.exists(input_path):
            os.remove(input_path)
        flash(f"❌ Erro ao converter PDF para Excel: {str(e)}")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: PDF para PowerPoint ---
@app.route('/pdf_para_ppt', methods=['POST'])
def pdf_para_ppt():
    try:
        # Verifica se o arquivo foi enviado no formulário
        if 'pdf' not in request.files:
            flash("⚠️ Nenhum arquivo foi enviado.")
            return redirect(request.referrer or url_for('index'))

        file = request.files['pdf']

        # Verifica se o usuário selecionou algum arquivo
        if file.filename == '':
            flash("⚠️ Nenhum arquivo selecionado.")
            return redirect(request.referrer or url_for('index'))

        # Verifica se o arquivo possui extensão .pdf (validação básica)
        if not file.filename.lower().endswith('.pdf'):
            flash(f"❌ O arquivo '{file.filename}' não possui extensão .pdf.")
            return redirect(request.referrer or url_for('index'))

        # Garante nome seguro e salva o PDF na pasta temporária
        filename = secure_filename(file.filename)
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        # Tenta abrir o PDF para garantir que não está protegido nem corrompido
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    os.remove(input_path)
                    flash("🔒 O arquivo PDF está protegido por senha.")
                    return redirect(request.referrer or url_for('index'))
        except PdfReadError:
            # PDF inválido, corrompido ou incompleto
            os.remove(input_path)
            flash("❌ O arquivo PDF parece estar incompleto ou corrompido.")
            return redirect(request.referrer or url_for('index'))

        # Define nome do PowerPoint de saída com timestamp legível
        output_filename = f"apresentacao_{datetime.now().strftime('%d-%m-%Y_%H-%M-%S')}.pptx"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Converte cada página do PDF em uma imagem (usado para inserir no PPT)
        imagens = convert_from_path(input_path)
        prs = Presentation()

        # Define tamanho padrão do slide em 16:9 (pode ser ajustado)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Para cada imagem gerada do PDF, cria um slide com a imagem ocupando todo o slide
        for img in imagens:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout em branco

            # Gera um caminho temporário para salvar a imagem JPEG (usa timestamp para evitar conflitos)
            img_path = os.path.join(RESULT_FOLDER, f"slide_temp_{datetime.now().timestamp()}.jpg")
            img.save(img_path, 'JPEG')

            # Adiciona a imagem ao slide, ajustando para preencher o slide inteiro
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

            # Remove a imagem temporária após adicioná-la para evitar acúmulo (opcional)
            os.remove(img_path)

        # Salva o arquivo PPTX final
        prs.save(output_path)

        # Remove o PDF temporário após a conversão
        os.remove(input_path)

        # Mensagem de sucesso e renderiza página de download
        flash("✅ PDF convertido para PowerPoint com sucesso!")
        return render_template('download_unico.html', filename=output_filename)

    except Exception as e:
        # Remove arquivo temporário em caso de erro
        if 'input_path' in locals() and os.path.exists(input_path):
            os.remove(input_path)
        flash(f"❌ Erro ao converter PDF para PowerPoint: {str(e)}")
        return redirect(request.referrer or url_for('index'))

# -_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-_-

# --- ROTA: Marca d'água no PDF ---
@app.route('/marca_dagua', methods=['POST'])
def marca_dagua():
    try:
        # Verifica se o arquivo PDF foi enviado pelo formulário
        if 'pdf' not in request.files:
            flash("⚠️ Nenhum arquivo foi enviado.")
            return redirect(request.referrer or url_for('index'))

        file = request.files['pdf']
        texto = request.form.get('texto', '').strip()

        # Valida se o arquivo possui nome e extensão .pdf
        if not file or file.filename == '' or not file.filename.lower().endswith('.pdf'):
            flash("❌ Envie um arquivo PDF válido.")
            return redirect(request.referrer or url_for('index'))

        # Valida se o texto da marca d'água foi fornecido
        if not texto:
            flash("❌ Insira um texto para a marca d'água.")
            return redirect(request.referrer or url_for('index'))

        # Salva o PDF recebido com nome seguro
        filename = secure_filename(file.filename)
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)

        # Verifica se o PDF está protegido ou corrompido
        try:
            with open(input_path, 'rb') as f:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    os.remove(input_path)
                    flash("🔒 O arquivo PDF está protegido por senha.")
                    return redirect(request.referrer or url_for('index'))
        except PdfReadError:
            os.remove(input_path)
            flash("❌ O arquivo PDF parece estar incompleto ou corrompido.")
            return redirect(request.referrer or url_for('index'))

        # Cria um PDF temporário para a marca d'água com ReportLab
        watermark_pdf = BytesIO()
        c = canvas.Canvas(watermark_pdf, pagesize=letter)
        c.setFont("Helvetica", 40)
        c.setFillGray(0.5, 0.3)  # Cor cinza translúcida para a marca
        c.saveState()
        c.translate(300, 400)
        c.rotate(45)
        c.drawCentredString(0, 0, texto)
        c.restoreState()
        c.save()
        watermark_pdf.seek(0)

        # Lê o PDF da marca d'água e o PDF original
        watermark = PdfReader(watermark_pdf)
        reader = PdfReader(input_path)
        writer = PdfWriter()

        watermark_page = watermark.pages[0]

        # Aplica a marca d'água em cada página do PDF original
        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)

        # Define o nome do arquivo de saída com timestamp legível
        output_filename = f"marcado_{datetime.now().strftime('%d-%m-%Y_%H-%M-%S')}.pdf"
        output_path = os.path.join(RESULT_FOLDER, output_filename)

        # Salva o novo PDF com marca d'água
        with open(output_path, 'wb') as f_out:
            writer.write(f_out)

        # Remove o PDF temporário de entrada para limpeza
        os.remove(input_path)

        flash("✅ Marca d'água adicionada com sucesso!")
        return render_template('download_unico.html', filename=output_filename)

    except Exception as e:
        # Remove o arquivo temporário em caso de erro
        if 'input_path' in locals() and os.path.exists(input_path):
            os.remove(input_path)
        flash(f"❌ Erro ao adicionar marca d'água: {str(e)}")
        return redirect(request.referrer or url_for('index'))

if __name__ == '__main__':
    app.run(debug=False)

